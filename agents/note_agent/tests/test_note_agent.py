from __future__ import annotations

from dataclasses import dataclass, field
from importlib.util import module_from_spec, spec_from_file_location
from io import BytesIO
import json
from pathlib import Path
from typing import Any

from fastapi.testclient import TestClient
from docx import Document
import openpyxl
from pptx import Presentation

from agents.note_agent.agent import NoteAgent
from agents.note_agent import agent as note_agent_module
from agents.note_agent.document_extractors import extract_document_bytes
from agents.note_agent.pdf_understanding import assess_pdf_text_quality, select_escalated_pages, select_initial_pages
from agents.note_agent.schemas import HealthStatus, IngestClassification, NoteFormattingPlan, NoteIngestRequest, NoteInvokeRequest

_APP_MAIN_PATH = Path("/app/app/main.py")
if not _APP_MAIN_PATH.exists():
    _APP_MAIN_PATH = Path(__file__).resolve().parents[3] / "apps" / "note-agent" / "app" / "main.py"
_APP_SPEC = spec_from_file_location("note_agent_app_main", _APP_MAIN_PATH)
assert _APP_SPEC and _APP_SPEC.loader
_APP_MODULE = module_from_spec(_APP_SPEC)
_APP_SPEC.loader.exec_module(_APP_MODULE)
app = _APP_MODULE.app


@dataclass
class _FakeAi:
    plan_result: NoteFormattingPlan
    calls: list[dict[str, Any]] = field(default_factory=list)

    def plan(self, *, message: str, attachments: list[dict[str, Any]] | None = None, metadata: dict[str, Any] | None = None) -> NoteFormattingPlan:
        self.calls.append({"message": message, "attachments": attachments or [], "metadata": metadata or {}})
        return self.plan_result

    def plan_ingested_source(self, context):
        self.calls.append({"context": context.model_dump()})
        return IngestClassification(
            title=self.plan_result.title,
            canonical_title=self.plan_result.canonical_title or self.plan_result.title,
            summary=self.plan_result.summary,
            details=self.plan_result.details,
            action_items=self.plan_result.action_items,
            tags=self.plan_result.tags,
            destination=self.plan_result.destination,
            collection_path=self.plan_result.collection_path,
            confidence=self.plan_result.confidence,
            note_kind=self.plan_result.note_kind,
            media_class="documents",
            source_date=self.plan_result.source_date,
            followups=self.plan_result.followups,
            classification_method="text",
            evidence_summary="Text extraction used for classification.",
        )

    def plan_scanned_pdf(self, *, context, page_images, page_numbers, extracted_text):
        self.calls.append({"scanned_pdf": {"context": context.model_dump(), "page_numbers": page_numbers, "page_count": len(page_images), "extracted_text": extracted_text}})
        return IngestClassification(
            title=self.plan_result.title,
            canonical_title=self.plan_result.canonical_title or self.plan_result.title,
            summary=self.plan_result.summary,
            details=self.plan_result.details,
            action_items=self.plan_result.action_items,
            tags=self.plan_result.tags,
            destination=self.plan_result.destination,
            collection_path=self.plan_result.collection_path,
            confidence=self.plan_result.confidence,
            note_kind=self.plan_result.note_kind,
            media_class="documents",
            source_date=self.plan_result.source_date,
            followups=self.plan_result.followups,
            classification_method="vision",
            evidence_summary="Summary derived from page-image analysis because extracted PDF text was insufficient.",
        )


@dataclass
class _FakeTools:
    created_notes: list[dict[str, Any]] = field(default_factory=list)
    moved: list[tuple[str, str]] = field(default_factory=list)
    inbox_files: list[dict[str, Any]] = field(default_factory=list)
    ready_inbox_files: list[dict[str, Any]] = field(default_factory=list)
    file_contents: dict[str, dict[str, Any]] = field(default_factory=dict)

    def discover_tools(self):
        return []

    def healthcheck(self) -> HealthStatus:
        return HealthStatus(ok=True, mcp_reachable=True, tools_discovered=["nc_webdav_write_file"])

    def create_note_in_inbox(self, *, title: str | None, content: str, template: str | None = None, session_id: str | None = None, destination: str = "Inbox") -> dict[str, Any]:
        payload = {"title": title, "content": content, "session_id": session_id, "destination": destination}
        self.created_notes.append(payload)
        return {"path": f"/Notes/FamilyCloud/{destination}/{session_id}-test-note.md", "appended": False}

    def write_markdown_note(self, path: str, content: str) -> dict[str, Any]:
        self.created_notes.append({"path": path, "content": content, "session_id": None, "destination": path})
        self.file_contents[path] = {"content": content, "content_type": "text/markdown"}
        return {"path": path, "result": {"ok": True}}

    def upload_media(self, raw_bytes: bytes, filename: str, destination: str | None = None, content_type: str | None = None) -> dict[str, Any]:
        return {"path": f"{destination}/{filename}", "result": {"ok": True}}

    def list_inbox_files(self) -> list[dict[str, Any]]:
        return list(self.inbox_files)

    def list_ready_inbox_files(self, scope: str | None = None, tag_name: str | None = None, limit: int | None = None) -> list[dict[str, Any]]:
        files = list(self.ready_inbox_files)
        return files if limit is None else files[:limit]

    def read(self, path: str) -> dict[str, Any]:
        return self.file_contents[path]

    def read_raw(self, path: str) -> dict[str, Any]:
        payload = dict(self.file_contents[path])
        if payload.get("encoding") == "base64":
            return payload
        raise KeyError(path)

    def ensure_directory(self, path: str) -> str:
        return path

    def path_exists(self, path: str) -> bool:
        return path in self.file_contents

    def ensure_unique_path(self, path: str) -> str:
        if path not in self.file_contents:
            return path
        stem, ext = path.rsplit(".", 1)
        index = 2
        while f"{stem}-{index}.{ext}" in self.file_contents:
            index += 1
        return f"{stem}-{index}.{ext}"

    def move(self, path: str, destination_path: str) -> dict[str, Any]:
        self.moved.append((path, destination_path))
        if path in self.file_contents:
            self.file_contents[destination_path] = self.file_contents.pop(path)
        return {"status_code": 201}


def test_note_agent_formats_and_files_note():
    tools = _FakeTools()
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Kitchen remodel",
            summary="Met with contractor.",
            details="Need estimate and permit check.",
            action_items=["Get estimate", "Check permits"],
            tags=["projects", "kitchen"],
            destination="Projects",
            confidence=0.9,
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)
    response = agent.run(
        NoteInvokeRequest(
            session_id="notes-1",
            message="Met with contractor. Need estimate for kitchen remodel. Also check permits.",
            actor="u@example.com",
            family_id=2,
        )
    )
    assert response.status == "ok"
    assert any("/Archive/Raw/" in item.path for item in response.created_items)
    assert any("/Projects/" in item.path and "kitchen-remodel" in item.path for item in response.created_items)
    assert "/Inbox/" in tools.created_notes[0]["path"]
    assert "## Original Message" in tools.created_notes[0]["content"]
    assert "/Projects/" in tools.created_notes[1]["path"]


def test_note_app_smoke(monkeypatch):
    tools = _FakeTools()
    agent = NoteAgent(
        ai=_FakeAi(
            NoteFormattingPlan(
                title="Captured note",
                summary="Short summary",
                details="Full details",
                action_items=[],
                tags=["inbox"],
                destination="Inbox",
                confidence=0.4,
            )
        ),
        tools=tools,
    )

    monkeypatch.setattr(_APP_MODULE, "note_tools", lambda: tools)
    monkeypatch.setattr(_APP_MODULE, "get_note_tools", lambda: tools)
    monkeypatch.setattr(_APP_MODULE, "get_note_agent", lambda: agent)

    client = TestClient(app)
    health = client.get("/healthz")
    assert health.status_code == 200
    assert health.json()["ok"] is True

    invoke = client.post(
        "/v1/agents/note/invoke",
        headers={"X-Dev-User": "u@example.com"},
        json={
            "session_id": "smoke-1",
            "message": "Quick note for the inbox",
            "actor": "u@example.com",
            "family_id": 2,
            "attachments": [],
        },
    )
    assert invoke.status_code == 200
    body = invoke.json()
    assert body["status"] == "ok"
    assert body["created_items"][0]["path"]


def test_invoke_records_raw_source_before_filing_polished_note():
    tools = _FakeTools()
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Family meeting follow-up",
            summary="Summarized follow-up items.",
            details="Need to review the budget and schedule dentist appointments.",
            action_items=["Review the budget", "Schedule dentist appointments"],
            tags=["family"],
            destination="Areas",
            confidence=0.92,
            source_date="2026-03-01",
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)

    response = agent.run(
        NoteInvokeRequest(
            session_id="family-1",
            message="Family meeting notes: review the budget and schedule dentist appointments.",
            actor="u@example.com",
            family_id=2,
        )
    )

    assert response.status == "ok"
    assert any(item.path.startswith("/Notes/FamilyCloud/Inbox/") for item in response.created_items)
    assert any("/Archive/Raw/" in item.path for item in response.created_items)
    assert any("/Areas/" in item.path and "family-meeting-follow-up" in item.path for item in response.created_items)
    assert "## Original Message" in tools.created_notes[0]["content"]
    assert "Family meeting notes: review the budget" in tools.created_notes[0]["content"]
    assert "Archive/Raw" in tools.created_notes[1]["content"]


def test_ingest_only_processes_mcp_ready_inbox_text_files():
    tools = _FakeTools(
        ready_inbox_files=[
            {
                "path": "/Notes/FamilyCloud/Inbox/ready-note.md",
                "name": "ready-note.md",
                "content_type": "text/markdown",
                "last_modified": "Sat, 28 Feb 2026 12:00:00 GMT",
            },
        ],
        file_contents={
            "/Notes/FamilyCloud/Inbox/ready-note.md": {
                "content": "# Contractor note\n\nNeed estimate and permit check.",
                "content_type": "text/markdown",
            },
        },
    )
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Contractor follow-up",
            summary="Need estimate.",
            details="Permit check too.",
            action_items=["Request estimate"],
            tags=["projects"],
            destination="Projects",
            confidence=0.9,
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)
    response = agent.ingest(NoteIngestRequest(actor="u@example.com", family_id=2, session_id="ingest-1"))
    assert response.status == "ok"
    assert response.processed_count == 1
    assert response.skipped_count == 0
    assert any("/Archive/Raw/" in item.path for item in response.created_items)
    assert any("/Projects/" in item.path and "contractor-follow-up" in item.path for item in response.created_items)
    assert tools.moved[0][0] == "/Notes/FamilyCloud/Inbox/ready-note.md"
    assert "/Projects/" in tools.created_notes[0]["path"]
    assert "contractor-follow-up" in tools.created_notes[0]["path"]


def test_ingest_ignores_filename_or_content_ready_heuristics_without_mcp_tag():
    tools = _FakeTools(
        inbox_files=[
            {
                "path": "/Notes/FamilyCloud/Inbox/ready-note.md",
                "name": "ready-note.md",
                "content_type": "text/markdown",
                "last_modified": "Sat, 28 Feb 2026 12:00:00 GMT",
            }
        ],
        ready_inbox_files=[],
        file_contents={
            "/Notes/FamilyCloud/Inbox/ready-note.md": {
                "content": "# Draft\n\nTags: ready\n\nStill drafting.",
                "content_type": "text/markdown",
            }
        },
    )
    agent = NoteAgent(
        ai=_FakeAi(
            NoteFormattingPlan(
                title="Should not run",
                summary="",
                details="",
                action_items=[],
                tags=["inbox"],
                destination="Inbox",
                confidence=0.9,
            )
        ),
        tools=tools,
    )

    response = agent.ingest(NoteIngestRequest(actor="u@example.com", family_id=2, session_id="ingest-2"))
    assert response.status == "ok"
    assert response.processed_count == 0
    assert not tools.created_notes


def test_media_ingest_uses_extracted_document_text_for_planning():
    tools = _FakeTools(
        ready_inbox_files=[
            {
                "path": "/Notes/FamilyCloud/Inbox/church-budget.docx",
                "name": "church-budget.docx",
                "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "last_modified": "Sat, 28 Feb 2026 12:00:00 GMT",
            }
        ],
        file_contents={
            "/Notes/FamilyCloud/Inbox/church-budget.docx": {
                "content": "Church budget review for March council meeting. Discuss missions, outreach, and youth ministry funding.",
                "content_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "parsed": True,
            }
        },
    )
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Church budget review",
            summary="Budget notes for council.",
            details="Church council budget planning.",
            action_items=["Review outreach budget"],
            tags=["church", "budget"],
            destination="Areas",
            confidence=0.88,
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)

    response = agent.ingest(NoteIngestRequest(actor="u@example.com", family_id=2, session_id="ingest-3"))
    assert response.status == "ok"
    assert response.processed_count == 1
    context = ai.calls[0]["context"]
    assert "Church budget review for March council meeting" in context["ocr_text"]
    assert tools.created_notes[0]["path"].startswith("/Notes/FamilyCloud/Areas")
    assert "Archive/Raw" in tools.created_notes[0]["content"]
    assert "Extraction mode" in tools.created_notes[0]["content"]


def test_polished_ingest_note_uses_nextcloud_files_app_link(monkeypatch):
    monkeypatch.setattr(note_agent_module.note_settings, "nextcloud_base_url", "https://nextcloud.family.callender")
    tools = _FakeTools(
        ready_inbox_files=[
            {
                "path": "/Notes/FamilyCloud/Inbox/ideate-whiteboard-sketches.whiteboard",
                "name": "ideate-whiteboard-sketches.whiteboard",
                "content_type": "application/json",
                "last_modified": "Sun, 01 Mar 2026 12:00:00 GMT",
            }
        ],
        file_contents={
            "/Notes/FamilyCloud/Inbox/ideate-whiteboard-sketches.whiteboard": {
                "content": '{"elements":[{"type":"text","text":"Ideate whiteboard sketches"}]}',
                "content_type": "application/json",
                "parsed": True,
            }
        },
    )
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Ideate whiteboard sketches",
            canonical_title="ideate whiteboard sketches",
            summary="Whiteboard sketch notes.",
            details="Initial ideation captured on whiteboard.",
            action_items=[],
            tags=["whiteboard"],
            destination="Archive",
            collection_path="Raw",
            confidence=0.9,
            source_date="2026-03-01",
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)

    response = agent.ingest(NoteIngestRequest(actor="u@example.com", family_id=2, session_id="ingest-4"))

    assert response.status == "ok"
    assert "https://nextcloud.family.callender/apps/files/?dir=/Notes/FamilyCloud/Archive/Raw/Raw/2026&relPath=2026-03-01-ideate-whiteboard-sketches-raw.whiteboard" in tools.created_notes[0]["content"]
    assert "- Extraction mode: `ocr`" in tools.created_notes[0]["content"]


def test_pdf_quality_gate_and_page_selection():
    weak = assess_pdf_text_quality("[image placeholder]\n[image placeholder]", page_count=13)
    weak_page_refs = assess_pdf_text_quality("pdf-0 pdf-1 pdf-2 pdf-3 pdf-4 pdf-5 pdf-6 pdf-7 pdf-8 pdf-9 pdf-10 pdf-11 pdf-12 page images are referenced", page_count=13)
    strong = assess_pdf_text_quality("This is a readable PDF with enough text to summarize." * 10, page_count=2)

    assert weak.quality == "image_only"
    assert weak_page_refs.quality == "image_only"
    assert strong.quality == "usable"
    assert select_initial_pages(13) == [1, 2, 7, 13]
    assert select_escalated_pages(13, initial_pages=[1, 2, 7, 13]) == [1, 2, 3, 4, 5, 7, 13]


def test_scanned_pdf_uses_vision_fallback(monkeypatch):
    monkeypatch.setattr(note_agent_module, "pdf_page_count", lambda raw: 13)
    monkeypatch.setattr(
        note_agent_module,
        "render_pdf_pages",
        lambda raw, pages: [{"page": page, "mime_type": "image/png", "base64": "ZmFrZQ=="} for page in pages],
    )
    tools = _FakeTools(
        ready_inbox_files=[
            {
                "path": "/Notes/FamilyCloud/Inbox/img20250919_19320220.pdf",
                "name": "img20250919_19320220.pdf",
                "content_type": "application/pdf",
                "last_modified": "Fri, 19 Sep 2025 19:32:20 GMT",
            }
        ],
        file_contents={
            "/Notes/FamilyCloud/Inbox/img20250919_19320220.pdf": {
                "content": "pdf-0 pdf-1 pdf-2 pdf-3 pdf-4 pdf-5 pdf-6 pdf-7 pdf-8 pdf-9 pdf-10 pdf-11 pdf-12 page images are referenced",
                "content_type": "application/pdf",
                "parsed": True,
                "bytes_base64": "JVBERi0xLjQKJUZBS0UK",
            }
        },
    )
    ai = _FakeAi(
        NoteFormattingPlan(
            title="Handwritten church notes",
            canonical_title="handwritten church notes",
            summary="Summary derived from handwritten church notes.",
            details="Likely sermon notes with handwritten references.",
            action_items=[],
            tags=["church", "scan"],
            destination="Areas",
            collection_path="Church",
            confidence=0.83,
            source_date="2025-09-19",
        )
    )
    agent = NoteAgent(ai=ai, tools=tools)

    response = agent.ingest(NoteIngestRequest(actor="u@example.com", family_id=2, session_id="ingest-5"))

    assert response.status == "ok"
    assert response.processed_count == 1
    scanned_call = ai.calls[0]["scanned_pdf"]
    assert scanned_call["page_numbers"] == [1, 2, 7, 13]
    assert "/Areas/Church/" in tools.created_notes[0]["path"]
    assert "- OCR quality: `image_only`" in tools.created_notes[0]["content"]
    assert "- Analyzed pages: `1, 2, 7, 13`" in tools.created_notes[0]["content"]
    assert "- Classification method: `hybrid`" in tools.created_notes[0]["content"]


def test_document_extractors_support_docx_pptx_xlsx_and_whiteboard():
    doc_buffer = BytesIO()
    doc = Document()
    doc.add_heading("Sunday Service 11-30-25", level=1)
    doc.add_paragraph("His name will be.")
    doc.save(doc_buffer)

    ppt_buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Wonderful Counselor"
    slide.placeholders[1].text = "Jesus knows exactly what we should do."
    presentation.save(ppt_buffer)

    xlsx_buffer = BytesIO()
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Budget"
    sheet.append(["Ministry", "Amount"])
    sheet.append(["Outreach", 500])
    workbook.save(xlsx_buffer)

    whiteboard_payload = {
        "elements": [
            {"type": "text", "text": "Launch small groups in January"},
            {"type": "rectangle", "label": "Prayer"},
        ]
    }

    assert "His name will be." in extract_document_bytes(doc_buffer.getvalue(), ".docx").text
    assert "Wonderful Counselor" in extract_document_bytes(ppt_buffer.getvalue(), ".pptx").text
    assert "Outreach | 500" in extract_document_bytes(xlsx_buffer.getvalue(), ".xlsx").text
    assert "Launch small groups in January" in extract_document_bytes(json.dumps(whiteboard_payload).encode("utf-8"), ".whiteboard").text


def test_note_app_ingest_smoke(monkeypatch):
    tools = _FakeTools(
        inbox_files=[],
    )
    agent = NoteAgent(
        ai=_FakeAi(
            NoteFormattingPlan(
                title="Captured note",
                summary="Short summary",
                details="Full details",
                action_items=[],
                tags=["inbox"],
                destination="Inbox",
                confidence=0.4,
            )
        ),
        tools=tools,
    )

    monkeypatch.setattr(_APP_MODULE, "note_tools", lambda: tools)
    monkeypatch.setattr(_APP_MODULE, "get_note_tools", lambda: tools)
    monkeypatch.setattr(_APP_MODULE, "get_note_agent", lambda: agent)

    client = TestClient(app)
    ingest = client.post(
        "/v1/agents/note/ingest",
        headers={"X-Dev-User": "u@example.com"},
        json={
            "session_id": "ingest-1",
            "actor": "u@example.com",
            "family_id": 2,
            "max_items": 10,
        },
    )
    assert ingest.status_code == 200
    body = ingest.json()
    assert body["status"] == "ok"
    assert body["processed_count"] == 0
