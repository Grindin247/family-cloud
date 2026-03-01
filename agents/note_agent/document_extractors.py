from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
import json
import os
import subprocess
import tempfile
import zipfile
from typing import Callable
from xml.etree import ElementTree as ET

import openpyxl
from docx import Document
from odf import teletype
from odf.opendocument import load as odf_load
from odf.text import H, P
from pptx import Presentation
import xlrd


@dataclass
class ExtractedDocument:
    text: str
    parser: str


def _normalize_lines(lines: list[str]) -> str:
    cleaned = [line.strip() for line in lines if line and line.strip()]
    return "\n".join(cleaned).strip()


def _run_external(command: list[str], suffix: str, raw_bytes: bytes) -> str:
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as handle:
        handle.write(raw_bytes)
        temp_path = handle.name
    try:
        result = subprocess.run(command + [temp_path], check=True, capture_output=True, text=True)
        return result.stdout.strip()
    finally:
        try:
            os.unlink(temp_path)
        except FileNotFoundError:
            pass


def _extract_docx(raw_bytes: bytes) -> ExtractedDocument:
    document = Document(BytesIO(raw_bytes))
    lines: list[str] = []
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            lines.append(paragraph.text)
    for table in document.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                lines.append(" | ".join(row_values))
    return ExtractedDocument(text=_normalize_lines(lines), parser="python-docx")


def _extract_doc(raw_bytes: bytes) -> ExtractedDocument:
    for command in (["antiword"], ["catdoc", "-w"]):
        try:
            text = _run_external(command, ".doc", raw_bytes)
            if text:
                return ExtractedDocument(text=text, parser=command[0])
        except Exception:
            continue
    raise RuntimeError("No DOC extractor available")


def _extract_pptx(raw_bytes: bytes) -> ExtractedDocument:
    presentation = Presentation(BytesIO(raw_bytes))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                slide_lines.append(text.strip())
        if slide_lines:
            lines.append(f"Slide {index}")
            lines.extend(slide_lines)
    return ExtractedDocument(text=_normalize_lines(lines), parser="python-pptx")


def _extract_ppt(raw_bytes: bytes) -> ExtractedDocument:
    text = _run_external(["catppt"], ".ppt", raw_bytes)
    return ExtractedDocument(text=text, parser="catppt")


def _extract_xlsx(raw_bytes: bytes) -> ExtractedDocument:
    workbook = openpyxl.load_workbook(BytesIO(raw_bytes), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value not in (None, "")]
            if values:
                lines.append(" | ".join(values))
    return ExtractedDocument(text=_normalize_lines(lines), parser="openpyxl")


def _extract_xls(raw_bytes: bytes) -> ExtractedDocument:
    workbook = xlrd.open_workbook(file_contents=raw_bytes)
    lines: list[str] = []
    for sheet in workbook.sheets():
        lines.append(f"Sheet: {sheet.name}")
        for row_idx in range(sheet.nrows):
            values = [str(sheet.cell_value(row_idx, col_idx)).strip() for col_idx in range(sheet.ncols)]
            values = [value for value in values if value]
            if values:
                lines.append(" | ".join(values))
    return ExtractedDocument(text=_normalize_lines(lines), parser="xlrd")


def _extract_odg(raw_bytes: bytes) -> ExtractedDocument:
    with tempfile.NamedTemporaryFile(delete=False, suffix=".odg") as handle:
        handle.write(raw_bytes)
        temp_path = handle.name
    try:
        doc = odf_load(temp_path)
        lines: list[str] = []
        for element_type in (H, P):
            for element in doc.getElementsByType(element_type):
                text = teletype.extractText(element).strip()
                if text:
                    lines.append(text)
        if not lines:
            with zipfile.ZipFile(BytesIO(raw_bytes)) as archive:
                content = archive.read("content.xml")
            root = ET.fromstring(content)
            for node in root.iter():
                if node.text and node.text.strip():
                    lines.append(node.text.strip())
        return ExtractedDocument(text=_normalize_lines(lines), parser="odfpy")
    finally:
        try:
            os.unlink(temp_path)
        except FileNotFoundError:
            pass


def _extract_whiteboard(raw_bytes: bytes) -> ExtractedDocument:
    payload = json.loads(raw_bytes.decode("utf-8"))
    lines: list[str] = []
    elements = payload.get("elements", [])
    for element in elements:
        if not isinstance(element, dict):
            continue
        text = str(element.get("text") or element.get("label") or "").strip()
        if text:
            lines.append(text)
        elif element.get("type"):
            lines.append(f"{element['type']}")
    return ExtractedDocument(text=_normalize_lines(lines), parser="whiteboard-json")


_EXTRACTORS: dict[str, Callable[[bytes], ExtractedDocument]] = {
    ".docx": _extract_docx,
    ".doc": _extract_doc,
    ".pptx": _extract_pptx,
    ".ppt": _extract_ppt,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xls,
    ".odg": _extract_odg,
    ".whiteboard": _extract_whiteboard,
}


def can_extract_extension(extension: str) -> bool:
    return extension.lower() in _EXTRACTORS


def extract_document_bytes(raw_bytes: bytes, extension: str) -> ExtractedDocument:
    extractor = _EXTRACTORS[extension.lower()]
    return extractor(raw_bytes)
